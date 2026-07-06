import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Colors
    DARK_BG = RGBColor(11, 15, 26)       # Deep space dark
    CARD_BG = RGBColor(20, 26, 45)       # Translucent card background
    TEXT_WHITE = RGBColor(255, 255, 255) # Clear headers
    TEXT_MUTED = RGBColor(148, 163, 184) # Muted subtitle text
    TEXT_BODY = RGBColor(203, 213, 225)  # Soft gray body text
    
    CYAN = RGBColor(6, 182, 212)        # Electric cyan (accent 1)
    PURPLE = RGBColor(139, 92, 246)      # Neon purple (accent 2)
    GREEN = RGBColor(16, 185, 129)       # Leaflet green (accent 3)

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide (Cover)
    # -------------------------------------------------------------
    slide_layout = prs.slide_layouts[6] # blank
    slide1 = prs.slides.add_slide(slide_layout)
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

    # Add center glow circle representation
    circle = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.66), Inches(0.75), Inches(6.0), Inches(6.0))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(15, 20, 38)
    circle.line.color.rgb = PURPLE
    circle.line.width = Pt(1.5)

    # Title textbox
    txBox = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "METROPOLIS AI"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Trebuchet MS'
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

    # Subtitle textbox
    txBox_sub = slide1.shapes.add_textbox(Inches(1.0), Inches(3.6), Inches(11.33), Inches(1.0))
    tf_sub = txBox_sub.text_frame
    tf_sub.word_wrap = True
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "Smart City Decision Intelligence Command Center"
    p_sub.alignment = PP_ALIGN.CENTER
    p_sub.font.name = 'Trebuchet MS'
    p_sub.font.size = Pt(20)
    p_sub.font.color.rgb = CYAN

    # Project details at bottom
    txBox_bot = slide1.shapes.add_textbox(Inches(1.0), Inches(5.8), Inches(11.33), Inches(0.8))
    tf_bot = txBox_bot.text_frame
    tf_bot.word_wrap = True
    p_bot = tf_bot.paragraphs[0]
    p_bot.text = "Integrated IoT Telemetry • LSTM/XGBoost Forecasts • Leaflet GIS • Gemini RAG Audits"
    p_bot.alignment = PP_ALIGN.CENTER
    p_bot.font.name = 'Calibri'
    p_bot.font.size = Pt(14)
    p_bot.font.color.rgb = TEXT_MUTED

    # Helper function to create standard slide templates
    def add_standard_slide(title_text):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG

        # Title text
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8))
        tf_t = title_box.text_frame
        tf_t.word_wrap = True
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.name = 'Trebuchet MS'
        p_t.font.size = Pt(28)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_WHITE

        # Neon divider line
        divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.15), Inches(12.33), Inches(0.02))
        divider.fill.solid()
        divider.fill.fore_color.rgb = CYAN
        divider.line.fill.background()

        return slide

    # Helper to add structured cards
    def add_card(slide, x, y, width, height, card_title, bullets, accent_color=PURPLE):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = RGBColor(40, 50, 75)
        card.line.width = Pt(1)

        # Card Title
        title_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.15), Inches(width - 0.3), Inches(0.4))
        tf_c = title_box.text_frame
        tf_c.word_wrap = True
        p_c = tf_c.paragraphs[0]
        p_c.text = card_title
        p_c.font.name = 'Trebuchet MS'
        p_c.font.size = Pt(16)
        p_c.font.bold = True
        p_c.font.color.rgb = accent_color

        # Card Bullets
        content_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.65), Inches(width - 0.3), Inches(height - 0.8))
        tf_b = content_box.text_frame
        tf_b.word_wrap = True

        for i, text in enumerate(bullets):
            p_b = tf_b.add_paragraph() if i > 0 else tf_b.paragraphs[0]
            p_b.text = "•  " + text
            p_b.font.name = 'Calibri'
            p_b.font.size = Pt(13)
            p_b.font.color.rgb = TEXT_BODY
            p_b.space_after = Pt(6)

    # -------------------------------------------------------------
    # SLIDE 2: Project Vision & Core Motivation
    # -------------------------------------------------------------
    slide2 = add_standard_slide("Executive Summary & Project Vision")
    add_card(slide2, 0.5, 1.6, 5.9, 5.2, "The Operational Problem", [
        "Metropolitan systems often operate in isolated silos (traffic, power, water, environment), preventing unified response plans.",
        "Warning systems rarely couple predictive weather telemetry directly with autonomous field actuators.",
        "System administrators lack explainable, traceable rationale logs for automated decision sequences."
    ], CYAN)
    add_card(slide2, 6.93, 1.6, 5.9, 5.2, "The Metropolis AI Solution", [
        "Consolidates real-time IoT feeds, municipal registries, and Earth Engine satellite scans into a single pane of glass.",
        "Bridges predictive climate algorithms directly to autonomous Pub/Sub emergency triggers (monsoons -> sluice gate deployments).",
        "RAG-configured conversational assistant offering full lineage audits of queried databases and XAI logic trees."
    ], PURPLE)

    # -------------------------------------------------------------
    # SLIDE 3: System Architecture
    # -------------------------------------------------------------
    slide3 = add_standard_slide("Full Stack System Architecture")
    add_card(slide3, 0.5, 1.6, 3.8, 5.2, "1. Ingestion & Storage", [
        "IoT telemetry streaming for Air Quality (AQI), reservoirs, and power lines.",
        "Mocked GCP BigQuery datastores containing historical sensor logs.",
        "Real-time sensor fluctuation script simulating live operational loads."
    ], CYAN)
    add_card(slide3, 4.76, 1.6, 3.8, 5.2, "2. Analytics & Agent Logic", [
        "Explainable AI (XAI) feature importance ranking modifiers.",
        "ADK Multi-Agent cooperative dialogs syncing Traffic, Environmental, Health, and Utility agents.",
        "LSTM and XGBoost dual-model predictive regression graphs."
    ], PURPLE)
    add_card(slide3, 9.03, 1.6, 3.8, 5.2, "3. Map & Portal Frontend", [
        "Glassmorphism UI using Vanilla HTML5, CSS3, and ES6 Javascript.",
        "Leaflet GIS map rendering CartoDB Dark tiles with custom pulsing markers.",
        "Focus City center translations adjusting all data parameters live."
    ], GREEN)

    # -------------------------------------------------------------
    # SLIDE 4: Real-Time Overview Dashboard
    # -------------------------------------------------------------
    slide4 = add_standard_slide("Overview Dashboard & Analytics")
    add_card(slide4, 0.5, 1.6, 5.9, 5.2, "BigQuery Analytics (Looker Mock)", [
        "KPI telemetry cards displaying live metrics for the focused city.",
        "Interactive Chart.js line graphs modeling grid electrical load trends and water flow rates.",
        "Automatic severity styling triggers matching AQI safety indexes.",
        "Sector Risk status logs classifying severity alerts in real-time."
    ], CYAN)
    add_card(slide4, 6.93, 1.6, 5.9, 5.2, "Operational Direct Actions", [
        "Focus City Selector dropdown immediately translating active coordinate zones (Mumbai, Tokyo, Sydney, London, São Paulo, Cairo, SF).",
        "Executive Report modal compiling live metrics to generate official directive files.",
        "Simulated digital signature verification using secure hash codes."
    ], GREEN)

    # -------------------------------------------------------------
    # SLIDE 5: Predictive Sandbox & GIS Mappings
    # -------------------------------------------------------------
    slide5 = add_standard_slide("Predictive Sandbox & Temporal GIS")
    add_card(slide5, 0.5, 1.6, 5.9, 5.2, "Predictive Sandbox Engine", [
        "Dual-model forecasting comparisons showing LSTM (smooth nodes) vs XGBoost (jagged trees) predictions.",
        "Environment variables (Rainfall, Temp offsets, Holidays) modifying charts live.",
        "Explainable AI (XAI) bar charts plotting input feature weights based on user modifications."
    ], PURPLE)
    add_card(slide5, 6.93, 1.6, 5.9, 5.2, "Temporal Earth Engine GIS", [
        "Leaflet layers mapping Flood Hazard Risk (blue), Crop NDVI health (green), and Urban Heat Islands (red).",
        "Time-lapse slider (2020 - 2026) altering layer radii to model climate impact.",
        "Pulsing IoT markers linking to vision overlays and drainage workflows."
    ], CYAN)

    # -------------------------------------------------------------
    # SLIDE 6: Multimodal AI & Autonomous Workflows
    # -------------------------------------------------------------
    slide6 = add_standard_slide("Multimodal AI & Autonomous Pipelines")
    add_card(slide6, 0.5, 1.6, 5.9, 5.2, "Multimodal AI Inspector", [
        "Vision AI frame analyzer rendering overlay bounding boxes with confidence intervals on garbage and road damage photos.",
        "CCTV Video tracking displaying real-time pedestrian grids and Safety status alerts.",
        "Voice NLP transcriber typing out citizen calls and parsing location, category, and urgency tags."
    ], GREEN)
    add_card(slide6, 6.93, 1.6, 5.9, 5.2, "Autonomous Action Pipelines", [
        "Structured Pub/Sub processing chains showing workflow progress nodes.",
        "Cross-Tab alert triggers: moving the rain forecast slider past 40mm automatically fires emergency sluice gate opening routines.",
        "Direct visual log alerts pushed to system operators in real-time."
    ], PURPLE)

    # -------------------------------------------------------------
    # SLIDE 7: GCP Deployment & Open-Source Repository
    # -------------------------------------------------------------
    slide7 = add_standard_slide("Production Deployment & Open-Source Sync")
    add_card(slide7, 0.5, 1.6, 5.9, 5.2, "GCP Cloud Storage Hosting", [
        "Deployed as a static, public web service on Google Cloud Storage.",
        "Configured with Uniform Bucket-Level Access and Public Object Viewers.",
        "Hosted bucket location: gs://metropolis-ai-536438030825 (us-central1).",
        "Production URL: https://storage.googleapis.com/metropolis-ai-536438030825/index.html"
    ], CYAN)
    add_card(slide7, 6.93, 1.6, 5.9, 5.2, "GitHub Repository Sync", [
        "Repository pushed to: https://github.com/devarkarrenika-pixel/Metropolis-AI",
        "Fully tracked to main branch, containing custom gitignore setups.",
        "Includes architectural flowcharts, technology descriptions, and setup guidelines in a professional README.md."
    ], PURPLE)

    # Save
    prs.save("Metropolis_AI_Presentation.pptx")
    print("Presentation compiled successfully!")

if __name__ == '__main__':
    build_presentation()
